import hashlib
import json
import os
import sys
import tempfile
import zipfile
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from urllib.parse import urlparse

CACHE_DIR = os.path.join(tempfile.gettempdir(), "opencode-doc-conv")
os.makedirs(CACHE_DIR, exist_ok=True)

DOC_EXTS = {".docx", ".doc"}
XLS_EXTS = {".xlsx", ".xls"}
PPT_EXTS = {".pptx", ".ppt"}
SUPPORTED_EXTS = DOC_EXTS | XLS_EXTS | PPT_EXTS


def _strip_zone_identifier(source_path):
    """Remove the Mark-of-the-Web (Zone.Identifier) alternate data stream so
    Office does not open the file in Protected View (read-only)."""
    try:
        import ctypes

        ctypes.windll.kernel32.DeleteFileW(source_path + ":Zone.Identifier")
    except Exception:
        pass


_OFFICE_EXES = {"EXCEL.EXE", "WINWORD.EXE", "POWERPNT.EXE"}
_LOCK_FILE = os.path.join(CACHE_DIR, "office_pids.json")
_managed = set()

_COM_MODULES = ("win32com", "pythoncom", "psutil")
_INSTALL_HINT = "pip install -r backend/requirements.txt"


def _ensure_deps(additional=(), require_com=True):
    import importlib

    required = _COM_MODULES if require_com else ()
    for mod in tuple(required) + tuple(additional):
        try:
            importlib.import_module(mod)
        except ImportError as exc:
            raise RuntimeError(
                f"Missing Python dependency '{mod}'. "
                f"Run: {_INSTALL_HINT}"
            ) from exc


def _office_pids():
    """Snapshot of the PIDs of running Office applications."""
    out = {}
    try:
        import psutil

        for proc in psutil.process_iter(["pid", "name"]):
            name = (proc.info["name"] or "").upper()
            if name in _OFFICE_EXES:
                out.setdefault(name, set()).add(proc.info["pid"])
    except Exception:
        pass
    return out


def _persist():
    try:
        with open(_LOCK_FILE, "w", encoding="utf-8") as f:
            json.dump({"managed": sorted([list(p) for p in _managed])}, f)
    except Exception:
        pass


def _reap(name, pid):
    try:
        import psutil

        proc = psutil.Process(pid)
        if proc.is_running() and (proc.name() or "").upper() == name:
            proc.kill()
            try:
                proc.wait(timeout=3)
            except psutil.TimeoutExpired:
                pass
    except (psutil.NoSuchProcess, psutil.AccessDenied):
        pass
    except Exception:
        pass


def _reap_stale():
    """Kill any Office processes we spawned in a previous (possibly crashed) run."""
    try:
        with open(_LOCK_FILE, "r", encoding="utf-8") as f:
            items = json.load(f).get("managed", [])
    except Exception:
        return
    for name, pid in items:
        _reap(name, pid)
    _managed.clear()
    _persist()


def _kill_new_office(before):
    """Force-quit any Office process spawned since the snapshot, so files
    opened by COM automation do not stay locked. Persists PIDs to disk so a
    crash cannot leave a zombie that holds a file lock."""
    for exe, old_pids in before.items():
        fresh = _office_pids().get(exe, set())
        for pid in fresh - old_pids:
            _managed.add((exe, pid))
            _persist()
            _reap(exe, pid)
            _managed.discard((exe, pid))
    _persist()


def cache_path(source_path):
    try:
        st = os.stat(source_path)
        key = f"{st.st_mtime_ns}:{st.st_size}:{os.path.abspath(source_path)}"
    except OSError:
        key = os.path.abspath(source_path)
    digest = hashlib.sha256(key.encode("utf-8")).hexdigest()
    return os.path.join(CACHE_DIR, digest + ".pdf")


def _export(source_path, tmp_path):
    import win32com.client

    ext = os.path.splitext(source_path)[1].lower()
    if ext in DOC_EXTS:
        app = win32com.client.DispatchEx("Word.Application")
        try:
            app.Visible = False
            app.DisplayAlerts = False
            doc = app.Documents.Open(source_path, ReadOnly=True)
            try:
                doc.SaveAs(tmp_path, FileFormat=17)
            finally:
                doc.Close(False)
        finally:
            app.Quit()
    elif ext in XLS_EXTS:
        app = win32com.client.DispatchEx("Excel.Application")
        try:
            app.Visible = False
            app.DisplayAlerts = False
            wb = app.Workbooks.Open(source_path, ReadOnly=True)
            try:
                for ws in wb.Worksheets:
                    ws.PageSetup.PrintHeadings = True
                    ws.PageSetup.PrintGridlines = True
                    ws.PageSetup.Zoom = False
                    ws.PageSetup.FitToPagesWide = 1
                    ws.PageSetup.FitToPagesTall = False
                wb.ExportAsFixedFormat(0, tmp_path)
            finally:
                wb.Close(False)
        finally:
            app.Quit()
    elif ext in PPT_EXTS:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(source_path, ReadOnly=True, WithWindow=False)
            try:
                pres.SaveAs(tmp_path, 32)
            finally:
                pres.Close()
        finally:
            app.Quit()


def convert(source_path, refresh=False):
    _strip_zone_identifier(source_path)
    ext = os.path.splitext(source_path)[1].lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Unsupported document type: {ext}")
    _ensure_deps()

    out_path = cache_path(source_path)
    if refresh:
        try:
            os.remove(out_path)
        except OSError:
            pass
    if os.path.exists(out_path):
        return out_path, True

    import pythoncom

    tmp_path = out_path[:-4] + ".conv.pdf"
    before = _office_pids()
    pythoncom.CoInitialize()
    try:
        _export(source_path, tmp_path)
    finally:
        pythoncom.CoUninitialize()
        _kill_new_office(before)

    if not os.path.exists(tmp_path):
        raise RuntimeError("Conversion produced no output")
    os.replace(tmp_path, out_path)
    return out_path, False


def _extract_pdf_text(source_path):
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader
    reader = PdfReader(source_path)
    pages = []
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _extract_word_text(source_path):
    import win32com.client

    app = win32com.client.DispatchEx("Word.Application")
    try:
        app.Visible = False
        app.DisplayAlerts = False
        doc = app.Documents.Open(source_path, ReadOnly=True)
        try:
            return doc.Content.Text
        finally:
            doc.Close(False)
    finally:
        app.Quit()


def _extract_excel_text(source_path):
    import win32com.client

    app = win32com.client.DispatchEx("Excel.Application")
    try:
        app.Visible = False
        app.DisplayAlerts = False
        wb = app.Workbooks.Open(source_path, ReadOnly=True)
        try:
            lines = []
            for ws in wb.Worksheets:
                lines.append(f"[Sheet: {ws.Name}]")
                used = ws.UsedRange
                if used is None:
                    continue
                rows = used.Rows.Count
                cols = used.Columns.Count
                for r in range(1, rows + 1):
                    values = []
                    for c in range(1, cols + 1):
                        value = None
                        try:
                            value = used.Cells(r, c).Value
                        except Exception:
                            pass
                        if value is not None:
                            values.append(str(value))
                    if any(v.strip() for v in values):
                        lines.append("\t".join(values))
            return "\n".join(lines)
        finally:
            wb.Close(False)
    finally:
        app.Quit()


def _extract_powerpoint_text(source_path):
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    try:
        pres = app.Presentations.Open(source_path, ReadOnly=True, WithWindow=False)
        try:
            lines = []
            for index, slide in enumerate(pres.Slides, 1):
                lines.append(f"[Slide {index}]")
                for shape in slide.Shapes:
                    if not shape.HasTextFrame:
                        continue
                    frame = shape.TextFrame
                    if not frame.HasText:
                        continue
                    lines.append(frame.TextRange.Text)
            return "\n".join(lines)
        finally:
            pres.Close()
    finally:
        app.Quit()


def _clean_msg_field(value):
    return (value or "").replace("\x00", "").strip()


def _extract_msg_text(source_path):
    import extract_msg

    msg = extract_msg.Message(source_path)
    try:
        lines = []
        sender = _clean_msg_field(str(getattr(msg, "sender", None) or ""))
        if sender:
            lines.append(f"From: {sender}")
        for label, attr in (("To", "to"), ("CC", "cc")):
            try:
                value = _clean_msg_field(getattr(msg, attr, None))
            except Exception:
                value = ""
            if value:
                lines.append(f"{label}: {value}")
        subject = _clean_msg_field(getattr(msg, "subject", None))
        if subject:
            lines.append(f"Subject: {subject}")
        try:
            date = _clean_msg_field(getattr(msg, "date", None))
        except Exception:
            date = ""
        if date:
            lines.append(f"Date: {date}")
        body = ""
        try:
            body = _clean_msg_field(getattr(msg, "body", None)) or body
        except Exception:
            body = ""
        if not body:
            html = ""
            try:
                html = _clean_msg_field(getattr(msg, "htmlBody", None)) or html
            except Exception:
                html = ""
            if html:
                import re

                body = re.sub(r"<[^>]+>", " ", html).strip()
                body = re.sub(r"\s+", " ", body)
        if body:
            lines.append("Body:")
            lines.append(body)
        return "\n".join(lines) or "(Empty email message)"
    finally:
        try:
            msg.close()
        except Exception:
            pass


def extract_text(source_path):
    _strip_zone_identifier(source_path)
    ext = os.path.splitext(source_path)[1].lower()
    if ext not in SUPPORTED_EXTS and ext not in {".pdf", ".msg"}:
        raise ValueError(f"Unsupported document type: {ext}")
    if ext == ".pdf":
        _ensure_deps(additional=("pypdf",), require_com=False)
        return _extract_pdf_text(source_path)
    if ext == ".msg":
        _ensure_deps(additional=("extract_msg",), require_com=False)
        return _extract_msg_text(source_path)
    _ensure_deps()

    import pythoncom

    before = _office_pids()
    pythoncom.CoInitialize()
    try:
        if ext in DOC_EXTS:
            return _extract_word_text(source_path)
        if ext in XLS_EXTS:
            return _extract_excel_text(source_path)
        if ext in PPT_EXTS:
            return _extract_powerpoint_text(source_path)
        raise ValueError(f"Unsupported document type: {ext}")
    finally:
        pythoncom.CoUninitialize()
        _kill_new_office(before)


def text_cache_path(source_path):
    base = cache_path(source_path)
    return base[:-4] + ".text.txt"


def _match_positions(concat, find, occurrence):
    positions = []
    start = 0
    while True:
        i = concat.find(find, start)
        if i == -1:
            break
        positions.append((i, i + len(find)))
        start = i + len(find)
    if occurrence:
        if 0 < occurrence <= len(positions):
            return [positions[occurrence - 1]]
        return []
    return positions


def _replace_in_runs(runs, find, replacement, occurrence):
    concat = "".join(r.text for r in runs)
    if find not in concat:
        return 0
    targets = _match_positions(concat, find, occurrence)
    if not targets:
        return 0
    pos = 0
    done = 0
    for r in runs:
        rstart = pos
        rend = rstart + len(r.text)
        segs = []
        cursor = rstart
        for (s, e) in targets:
            if s >= rend:
                break
            if e <= rstart:
                continue
            if s < rstart:
                cursor = e
                continue
            if s > cursor:
                segs.append(r.text[cursor - rstart:s - rstart])
            segs.append(replacement)
            cursor = e
            done += 1
        if cursor < rend:
            segs.append(r.text[cursor - rstart:])
        r.text = "".join(segs)
        pos = rend
    return done


def _docx_insert_paragraph_after(p, text):
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    new_p = p._p.makeelement(qn("w:p"), {})
    p._p.addnext(new_p)
    new_para = Paragraph(new_p, p._parent)
    new_para.add_run(text)


def _docx_all_paragraphs(doc):
    paras = list(doc.paragraphs)

    def walk(tables):
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    paras.extend(cell.paragraphs)
                    walk(cell.tables)

    walk(doc.tables)
    return paras


def _edit_docx(path, operations):
    import docx

    doc = docx.Document(path)
    paras = _docx_all_paragraphs(doc)
    results = []
    for op in operations:
        kind = op.get("op")
        applied = False
        if kind == "replace":
            find = op.get("find", "")
            if find:
                repl = op.get("replace", "")
                occurrence = op.get("occurrence") or 0
                budget = occurrence
                for p in paras:
                    if p.runs and find in p.text:
                        if occurrence:
                            made = _replace_in_runs(p.runs, find, repl, budget)
                            if made:
                                budget -= 1
                                applied = True
                            if budget <= 0:
                                break
                        else:
                            applied = _replace_in_runs(p.runs, find, repl, 0) > 0 or applied
        elif kind in ("insert_after", "insert_before"):
            find = op.get("find", "")
            if find:
                idx = op.get("occurrence") or 1
                for p in paras:
                    if find in p.text:
                        idx -= 1
                        if idx == 0:
                            if kind == "insert_after":
                                _docx_insert_paragraph_after(p, op.get("text", ""))
                            else:
                                p.insert_paragraph_before(op.get("text", ""))
                            applied = True
                            break
        elif kind == "append":
            doc.add_paragraph(op.get("text", ""))
            applied = True
        elif kind == "prepend":
            if doc.paragraphs:
                doc.paragraphs[0].insert_paragraph_before(op.get("text", ""))
                applied = True
        elif kind == "delete":
            find = op.get("find", "")
            if find:
                occurrence = op.get("occurrence") or 0
                budget = occurrence
                for p in paras:
                    if p.runs and find in p.text:
                        if occurrence:
                            made = _replace_in_runs(p.runs, find, "", budget)
                            if made:
                                budget -= 1
                                applied = True
                            if budget <= 0:
                                break
                        else:
                            applied = _replace_in_runs(p.runs, find, "", 0) > 0 or applied
        results.append({"op": kind, "applied": applied})
    doc.save(path)
    return results


def _cell_replace(ws, find, repl, occurrence, budget):
    for row in ws.iter_rows():
        for cell in row:
            if isinstance(cell.value, str) and find in cell.value:
                if occurrence:
                    if budget <= 0:
                        return budget
                    budget -= 1
                    cell.value = cell.value.replace(find, repl, 1)
                else:
                    cell.value = cell.value.replace(find, repl)
    return budget


def _edit_xlsx(path, operations):
    from openpyxl import load_workbook

    wb = load_workbook(path)
    results = []
    for op in operations:
        kind = op.get("op")
        applied = False
        if kind in ("replace", "delete"):
            find = op.get("find", "")
            if find:
                repl = "" if kind == "delete" else op.get("replace", "")
                occurrence = op.get("occurrence") or 0
                budget = occurrence
                for ws in wb.worksheets:
                    budget = _cell_replace(ws, find, repl, occurrence, budget)
                    if occurrence and budget <= 0:
                        break
                applied = True
        elif kind in ("insert_after", "insert_before"):
            find = op.get("find", "")
            if find:
                idx = op.get("occurrence") or 1
                for ws in wb.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            if isinstance(cell.value, str) and find in cell.value:
                                idx -= 1
                                if idx == 0:
                                    r = cell.row
                                    c = cell.column
                                    if kind == "insert_after":
                                        ws.insert_rows(r + 1)
                                        ws.cell(row=r + 1, column=c, value=op.get("text", ""))
                                    else:
                                        ws.insert_rows(r)
                                        ws.cell(row=r, column=c, value=op.get("text", ""))
                                    applied = True
                                    break
                        if applied:
                            break
                    if applied:
                        break
        elif kind == "append":
            wb.worksheets[0].append([op.get("text", "")])
            applied = True
        elif kind == "prepend":
            ws = wb.worksheets[0]
            ws.insert_rows(1)
            ws.cell(row=1, column=1, value=op.get("text", ""))
            applied = True
        results.append({"op": kind, "applied": applied})
    wb.save(path)
    return results


def _pptx_text_frames(prs):
    frames = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                frames.append(shape.text_frame)
    return frames


def _edit_pptx(path, operations):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(path)
    frames = _pptx_text_frames(prs)
    results = []
    for op in operations:
        kind = op.get("op")
        applied = False
        if kind in ("replace", "delete"):
            find = op.get("find", "")
            if find:
                repl = "" if kind == "delete" else op.get("replace", "")
                occurrence = op.get("occurrence") or 0
                budget = occurrence
                for tf in frames:
                    for para in tf.paragraphs:
                        runs = para.runs
                        if not runs or find not in "".join(r.text for r in runs):
                            continue
                        if occurrence:
                            made = _replace_in_runs(runs, find, repl, budget)
                            if made:
                                budget -= 1
                                applied = True
                            if budget <= 0:
                                break
                        else:
                            applied = _replace_in_runs(runs, find, repl, 0) > 0 or applied
                    if occurrence and budget <= 0:
                        break
        elif kind in ("insert_after", "insert_before"):
            find = op.get("find", "")
            if find:
                idx = op.get("occurrence") or 1
                for tf in frames:
                    paras = tf.paragraphs
                    for n, para in enumerate(paras):
                        if find in "".join(r.text for r in para.runs):
                            idx -= 1
                            if idx == 0:
                                if kind == "insert_after":
                                    if n + 1 < len(paras):
                                        paras[n + 1].insert_paragraph_before(op.get("text", ""))
                                    else:
                                        from pptx.oxml.ns import qn

                                        tf._txBody.append(tf._txBody.makeelement(qn("a:p"), {}))
                                        tf.paragraphs[-1].text = op.get("text", "")
                                else:
                                    para.insert_paragraph_before(op.get("text", ""))
                                applied = True
                                break
                    if applied:
                        break
        elif kind == "append":
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
            box.text_frame.text = op.get("text", "")
            applied = True
        elif kind == "prepend":
            if frames and frames[0].paragraphs:
                frames[0].paragraphs[0].insert_paragraph_before(op.get("text", ""))
                applied = True
        results.append({"op": kind, "applied": applied})
    prs.save(path)
    return results


def _edit_word_com(path, operations):
    import win32com.client

    app = win32com.client.DispatchEx("Word.Application")
    app.Visible = False
    app.DisplayAlerts = False
    doc = app.Documents.Open(path)
    results = []
    try:
        for op in operations:
            kind = op.get("op")
            applied = False
            if kind in ("replace", "delete"):
                find = op.get("find", "")
                if find:
                    repl = "" if kind == "delete" else op.get("replace", "")
                    occurrence = op.get("occurrence") or 0
                    rng = doc.Content
                    f = rng.Find
                    f.ClearFormatting()
                    f.Replacement.ClearFormatting()
                    if occurrence == 0:
                        if f.Execute(FindText=find, ReplaceWith=repl, Replace=2):
                            applied = True
                    else:
                        count = 0
                        f.Text = find
                        f.Forward = True
                        f.Wrap = 0
                        while f.Execute():
                            count += 1
                            if count == occurrence:
                                rng.Text = repl
                                applied = True
                                break
                            rng.Collapse(1)
            elif kind in ("insert_after", "insert_before"):
                find = op.get("find", "")
                if find:
                    idx = op.get("occurrence") or 1
                    rng = doc.Content
                    f = rng.Find
                    f.ClearFormatting()
                    f.Text = find
                    f.Forward = True
                    f.Wrap = 0
                    count = 0
                    while f.Execute():
                        count += 1
                        if count == idx:
                            if kind == "insert_after":
                                rng.InsertAfter(op.get("text", ""))
                            else:
                                rng.InsertBefore(op.get("text", ""))
                            applied = True
                            break
                        rng.Collapse(1)
            elif kind == "append":
                rng = doc.Content
                rng.Collapse(1)
                rng.InsertAfter("\r\n" + op.get("text", ""))
                applied = True
            elif kind == "prepend":
                rng = doc.Content
                rng.Collapse(0)
                rng.InsertBefore(op.get("text", ""))
                applied = True
            results.append({"op": kind, "applied": applied})
        doc.Save()
    finally:
        doc.Close(False)
        app.Quit()
    return results


def _edit_excel_com(path, operations):
    import win32com.client

    app = win32com.client.DispatchEx("Excel.Application")
    app.Visible = False
    app.DisplayAlerts = False
    wb = app.Workbooks.Open(path)
    results = []
    try:
        for op in operations:
            kind = op.get("op")
            applied = False
            if kind in ("replace", "delete"):
                find = op.get("find", "")
                if find:
                    repl = "" if kind == "delete" else op.get("replace", "")
                    occurrence = op.get("occurrence") or 0
                    budget = occurrence
                    for ws in wb.Worksheets:
                        used = ws.UsedRange
                        if used is None:
                            continue
                        rows = used.Rows.Count
                        cols = used.Columns.Count
                        for r in range(1, rows + 1):
                            for c in range(1, cols + 1):
                                cell = used.Cells(r, c)
                                value = None
                                try:
                                    value = cell.Value
                                except Exception:
                                    pass
                                if isinstance(value, str) and find in value:
                                    if occurrence:
                                        if budget <= 0:
                                            break
                                        budget -= 1
                                        cell.Value = value.replace(find, repl, 1)
                                    else:
                                        cell.Value = value.replace(find, repl)
                                    applied = True
                            if occurrence and budget <= 0:
                                break
                        if occurrence and budget <= 0:
                            break
            elif kind in ("insert_after", "insert_before"):
                find = op.get("find", "")
                if find:
                    idx = op.get("occurrence") or 1
                    for ws in wb.Worksheets:
                        used = ws.UsedRange
                        if used is None:
                            continue
                        rows = used.Rows.Count
                        cols = used.Columns.Count
                        found = False
                        for r in range(1, rows + 1):
                            for c in range(1, cols + 1):
                                value = None
                                try:
                                    value = used.Cells(r, c).Value
                                except Exception:
                                    pass
                                if isinstance(value, str) and find in value:
                                    idx -= 1
                                    if idx == 0:
                                        target = ws.Cells(r + 1 if kind == "insert_after" else r, 1)
                                        if kind == "insert_after":
                                            ws.Rows(r + 1).Insert()
                                        else:
                                            ws.Rows(r).Insert()
                                        ws.Cells(r + 1 if kind == "insert_after" else r, c).Value = op.get("text", "")
                                        applied = True
                                        found = True
                                        break
                            if found:
                                break
                        if found:
                            break
            elif kind == "append":
                ws = wb.Worksheets(1)
                last = ws.Cells(ws.Rows.Count, 1).End(-4162).Row
                ws.Cells(last + 1, 1).Value = op.get("text", "")
                applied = True
            elif kind == "prepend":
                ws = wb.Worksheets(1)
                ws.Rows(1).Insert()
                ws.Cells(1, 1).Value = op.get("text", "")
                applied = True
            results.append({"op": kind, "applied": applied})
        wb.Save()
    finally:
        wb.Close(False)
        app.Quit()
    return results


def _edit_powerpoint_com(path, operations):
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    pres = app.Presentations.Open(path, WithWindow=False)
    results = []
    try:
        for op in operations:
            kind = op.get("op")
            applied = False
            if kind in ("replace", "delete"):
                find = op.get("find", "")
                if find:
                    repl = "" if kind == "delete" else op.get("replace", "")
                    occurrence = op.get("occurrence") or 0
                    budget = occurrence
                    for slide in pres.Slides:
                        for shape in slide.Shapes:
                            if not shape.HasTextFrame:
                                continue
                            tf = shape.TextFrame
                            if not tf.HasText:
                                continue
                            tr = tf.TextRange
                            txt = tr.Text
                            if find not in txt:
                                continue
                            if occurrence:
                                if budget <= 0:
                                    break
                                n = txt.find(find)
                                count = 1
                                while n != -1 and count < budget:
                                    count += 1
                                    n = txt.find(find, n + len(find))
                                if n != -1:
                                    tr.Text = txt[:n] + repl + txt[n + len(find):]
                                    budget -= 1
                                    applied = True
                                if budget <= 0:
                                    break
                            else:
                                tr.Text = txt.replace(find, repl)
                                applied = True
                        if occurrence and budget <= 0:
                            break
                    if occurrence and budget <= 0:
                        break
            elif kind in ("insert_after", "insert_before"):
                find = op.get("find", "")
                if find:
                    idx = op.get("occurrence") or 1
                    done = False
                    for slide in pres.Slides:
                        for shape in slide.Shapes:
                            if not shape.HasTextFrame:
                                continue
                            tf = shape.TextFrame
                            if not tf.HasText:
                                continue
                            tr = tf.TextRange
                            txt = tr.Text
                            n = txt.find(find)
                            if n == -1:
                                continue
                            idx -= 1
                            if idx == 0:
                                pos = n + len(find) if kind == "insert_after" else n
                                tr.Text = txt[:pos] + op.get("text", "") + txt[pos:]
                                applied = True
                                done = True
                                break
                        if done:
                            break
            elif kind == "append":
                slide = pres.Slides.Add(pres.Slides.Count + 1, 12)
                box = slide.Shapes.AddTextbox(1, 20, 20, 700, 400)
                box.TextFrame.TextRange.Text = op.get("text", "")
                applied = True
            elif kind == "prepend":
                slide = pres.Slides(1)
                found = False
                for shape in slide.Shapes:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        tr = shape.TextFrame.TextRange
                        tr.Text = op.get("text", "") + tr.Text
                        applied = True
                        found = True
                        break
                if not found:
                    box = slide.Shapes.AddTextbox(1, 20, 20, 700, 400)
                    box.TextFrame.TextRange.Text = op.get("text", "")
                    applied = True
            results.append({"op": kind, "applied": applied})
        pres.Save()
    finally:
        pres.Close()
        app.Quit()
    return results


def _edit_legacy(path, ext, operations):
    if ext in DOC_EXTS:
        return _edit_word_com(path, operations)
    if ext in XLS_EXTS:
        return _edit_excel_com(path, operations)
    return _edit_powerpoint_com(path, operations)


def edit_document(source_path, operations):
    _strip_zone_identifier(source_path)
    ext = os.path.splitext(source_path)[1].lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Unsupported document type: {ext}")

    if ext in {".docx", ".xlsx", ".pptx"}:
        if zipfile.is_zipfile(source_path):
            _ensure_deps(additional=({"docx": "docx", "xlsx": "openpyxl", "pptx": "pptx"}[ext],), require_com=False)
            if ext == ".docx":
                return _edit_docx(source_path, operations)
            if ext == ".xlsx":
                return _edit_xlsx(source_path, operations)
            return _edit_pptx(source_path, operations)
        _ensure_deps()
        return _edit_legacy(source_path, ext, operations)

    _ensure_deps()
    import pythoncom

    before = _office_pids()
    pythoncom.CoInitialize()
    try:
        return _edit_legacy(source_path, ext, operations)
    finally:
        pythoncom.CoUninitialize()
        _kill_new_office(before)


class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        sys.stderr.write("[doc-converter] " + (fmt % args) + "\n")

    def _json(self, code, obj):
        body = json.dumps(obj).encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def do_GET(self):
        parsed = urlparse(self.path)
        if parsed.path == "/health":
            self._json(200, {"ok": True})
        else:
            self._json(404, {"error": "Not found"})

    def do_POST(self):
        parsed = urlparse(self.path)
        if parsed.path == "/extract":
            self._handle_extract()
            return
        if parsed.path == "/edit":
            self._handle_edit()
            return
        if parsed.path != "/convert":
            self._json(404, {"error": "Not found"})
            return
        try:
            length = int(self.headers.get("Content-Length") or 0)
            payload = json.loads(self.rfile.read(length) or b"{}")
            source_path = payload.get("path")
            if not source_path or not os.path.isfile(source_path):
                self._json(400, {"error": "Invalid path"})
                return
            refresh = bool(payload.get("refresh"))
            pdf_path, cached = convert(source_path, refresh)
            self._json(200, {"pdfPath": pdf_path, "cached": cached})
        except Exception as exc:
            self._json(500, {"error": str(exc)})

    def _handle_extract(self):
        try:
            length = int(self.headers.get("Content-Length") or 0)
            payload = json.loads(self.rfile.read(length) or b"{}")
            source_path = payload.get("path")
            if not source_path or not os.path.isfile(source_path):
                self._json(400, {"error": "Invalid path"})
                return
            out_path = text_cache_path(source_path)
            data = None
            if os.path.exists(out_path):
                if payload.get("refresh"):
                    try:
                        os.remove(out_path)
                    except OSError:
                        pass
                else:
                    try:
                        data = open(out_path, "r", encoding="utf-8").read()
                    except Exception:
                        data = None
            if data is None:
                data = extract_text(source_path)
                try:
                    with open(out_path, "w", encoding="utf-8") as f:
                        f.write(data)
                except Exception:
                    pass
            self._json(200, {"text": data, "fileName": os.path.basename(source_path)})
        except Exception as exc:
            self._json(500, {"error": str(exc)})

    def _handle_edit(self):
        try:
            length = int(self.headers.get("Content-Length") or 0)
            payload = json.loads(self.rfile.read(length) or b"{}")
            source_path = payload.get("path")
            operations = payload.get("operations")
            if not source_path or not os.path.isfile(source_path):
                self._json(400, {"error": "Invalid path"})
                return
            if not isinstance(operations, list) or not operations:
                self._json(400, {"error": "No operations provided"})
                return
            results = edit_document(source_path, operations)
            self._json(200, {"edited": True, "fileName": os.path.basename(source_path), "results": results})
        except Exception as exc:
            self._json(500, {"error": str(exc)})


def main():
    _reap_stale()
    port = int(os.environ.get("DOC_CONVERTER_PORT", "8765"))
    server = ThreadingHTTPServer(("127.0.0.1", port), Handler)
    server.serve_forever()


if __name__ == "__main__":
    main()
